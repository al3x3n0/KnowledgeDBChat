"""
Text processing service for document content extraction and chunking.
"""

import os
import re
from typing import List, Optional, Dict, Any, Tuple
import numpy as np
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.document_loaders import (
    TextLoader,
    PyPDFLoader,
    UnstructuredWordDocumentLoader,
    UnstructuredHTMLLoader,
    UnstructuredMarkdownLoader,
    UnstructuredPowerPointLoader
)
from sentence_transformers import SentenceTransformer
from loguru import logger
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

from app.core.config import settings


class TextProcessor:
    """Service for processing and extracting text from various document formats."""
    
    def __init__(self):
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
            separators=["\n\n", "\n", " ", ""]
        )
        # Initialize embedding model for semantic chunking if needed
        self.semantic_model = None
        if settings.RAG_CHUNKING_STRATEGY == "semantic":
            try:
                self.semantic_model = SentenceTransformer(settings.EMBEDDING_MODEL)
                logger.info(f"Initialized semantic chunking model: {settings.EMBEDDING_MODEL}")
            except Exception as e:
                logger.warning(f"Failed to initialize semantic chunking model: {e}. Falling back to fixed-size chunking.")
                self.semantic_model = None
    
    async def extract_text(self, file_path: str, content_type: Optional[str] = None) -> Tuple[str, Optional[Dict[str, Any]]]:
        """
        Extract text content from a file.
        
        Args:
            file_path: Path to the file to extract text from
            content_type: Optional MIME type of the file (e.g., 'application/pdf')
            
            Returns:
                Tuple of (text content, optional metadata dict)
            
        Raises:
            FileNotFoundError: If the file does not exist
            Exception: If text extraction fails for the file format
        """
        try:
            file_extension = os.path.splitext(file_path)[1].lower()
            metadata: Optional[Dict[str, Any]] = None
            
            # Determine the appropriate loader based on file extension or content type
            if file_extension == '.pdf' or (content_type and 'pdf' in content_type):
                text = await self._extract_pdf(file_path)
            elif file_extension in ['.docx', '.doc'] or (content_type and 'word' in content_type):
                text = await self._extract_word(file_path)
            elif file_extension in ['.pptx', '.ppt'] or (content_type and ('powerpoint' in content_type.lower() or 'presentation' in content_type.lower())):
                text, metadata = await self._extract_powerpoint(file_path)
            elif file_extension in ['.html', '.htm'] or (content_type and 'html' in content_type):
                text = await self._extract_html(file_path)
            elif file_extension in ['.md', '.markdown'] or (content_type and 'markdown' in content_type.lower()):
                text = await self._extract_markdown(file_path)
            elif file_extension == '.txt' or (content_type and 'text' in content_type):
                text = await self._extract_text_file(file_path)
            else:
                # Default to text extraction
                text = await self._extract_text_file(file_path)
            
            return text, metadata
                
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {e}")
            return "", None
    
    async def _extract_pdf(self, file_path: str) -> str:
        """Extract text from PDF file."""
        try:
            logger.info(f"Extracting text from PDF: {file_path}")
            loader = PyPDFLoader(file_path)
            documents = loader.load()
            
            if not documents:
                logger.warning(f"No content extracted from PDF: {file_path}")
                return ""
            
            # Combine all pages with page breaks
            text_content = "\n\n".join([doc.page_content for doc in documents])
            
            # Log extraction stats
            page_count = len(documents)
            text_length = len(text_content)
            logger.info(f"Extracted {text_length} characters from {page_count} pages of PDF: {file_path}")
            
            return text_content
        except ImportError as e:
            logger.error(f"PDF extraction library not available: {e}. Please install pypdf.")
            raise
        except Exception as e:
            logger.error(f"Error extracting PDF {file_path}: {e}", exc_info=True)
            return ""
    
    async def _extract_word(self, file_path: str) -> str:
        """Extract text from Word document."""
        try:
            loader = UnstructuredWordDocumentLoader(file_path)
            documents = loader.load()
            return "\n\n".join([doc.page_content for doc in documents])
        except Exception as e:
            logger.error(f"Error extracting Word document {file_path}: {e}")
            return ""
    
    async def _extract_powerpoint(self, file_path: str) -> Tuple[str, Optional[Dict[str, Any]]]:
        """Extract text and metadata from PowerPoint presentation, including notes/comments."""
        try:
            logger.info(f"Extracting text from PowerPoint: {file_path}")
            presentation = Presentation(file_path)
            slide_entries: List[Dict[str, Any]] = []
            slide_sections: List[str] = []
            
            for idx, slide in enumerate(presentation.slides, start=1):
                body_texts: List[str] = []
                title_text: Optional[str] = None
                
                for shape in slide.shapes:
                    text = self._get_shape_text(shape)
                    if not text:
                        continue
                    body_texts.append(text)
                    try:
                        if title_text is None and getattr(shape, "is_placeholder", False):
                            placeholder_type = shape.placeholder_format.type
                            if placeholder_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                                title_text = text.splitlines()[0]
                    except Exception:
                        pass
                
                slide_body = "\n".join(body_texts).strip()
                if not title_text and slide_body:
                    title_text = slide_body.splitlines()[0]
                
                notes_text = ""
                try:
                    if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                        note_paragraphs = []
                        for paragraph in slide.notes_slide.notes_text_frame.paragraphs:
                            note_line = "".join(run.text for run in paragraph.runs).strip()
                            if note_line:
                                note_paragraphs.append(note_line)
                        notes_text = "\n".join(note_paragraphs).strip()
                except Exception:
                    notes_text = ""
                
                comments: List[Dict[str, Any]] = []
                try:
                    comments_part = getattr(slide.part, "comments_part", None)
                    if comments_part is not None:
                        for comment in getattr(comments_part, "comment_list", []):
                            text = getattr(comment, "text", "") or ""
                            if not text.strip():
                                continue
                            author_name = None
                            try:
                                author_obj = getattr(comment, "author", None)
                                author_name = getattr(author_obj, "name", None) if author_obj else None
                                author_initials = getattr(author_obj, "initials", None) if author_obj else None
                            except Exception:
                                author_obj = getattr(comment, "author", None)
                                author_name = getattr(author_obj, "name", None) if author_obj else str(author_obj)
                                author_initials = None
                            timestamp = None
                            comment_dt = getattr(comment, "datetime", None) or getattr(comment, "dt", None)
                            if comment_dt:
                                try:
                                    timestamp = comment_dt.isoformat()
                                except Exception:
                                    timestamp = str(comment_dt)
                            comments.append({
                                "author": author_name,
                                "author_initials": author_initials,
                                "text": text.strip(),
                                "created_at": timestamp,
                            })
                except Exception:
                    pass
                
                slide_entry: Dict[str, Any] = {
                    "index": idx,
                }
                if title_text:
                    slide_entry["title"] = title_text
                if slide_body:
                    slide_entry["text"] = slide_body
                if notes_text:
                    slide_entry["notes"] = notes_text
                if comments:
                    slide_entry["comments"] = comments
                slide_entries.append(slide_entry)
                
                section_parts = [f"Slide {idx}: {title_text or 'Untitled slide'}"]
                if slide_body:
                    section_parts.append(slide_body)
                if notes_text:
                    section_parts.append(f"Notes:\n{notes_text}")
                if comments:
                    comment_lines = [
                        f"- {(c.get('author') or c.get('author_initials') or 'Comment')}: {c['text']}"
                        for c in comments
                    ]
                    section_parts.append("Comments:\n" + "\n".join(comment_lines))
                slide_sections.append("\n\n".join(part for part in section_parts if part))
            
            if slide_sections:
                text_content = "\n\n--- Slide ---\n\n".join(slide_sections)
                metadata = {
                    "presentation": {
                        "slide_count": len(slide_entries),
                        "slides": slide_entries,
                        "has_comments": any(entry.get("comments") for entry in slide_entries),
                    }
                }
                text_length = len(text_content)
                logger.info(f"Extracted {text_length} characters from {len(slide_entries)} slides of PowerPoint: {file_path}")
                return text_content, metadata
            
            # Fallback to unstructured loader if no sections were produced
            loader = UnstructuredPowerPointLoader(file_path)
            documents = loader.load()
            if not documents:
                logger.warning(f"No content extracted from PowerPoint: {file_path}")
                return "", None
            text_content = "\n\n--- Slide ---\n\n".join([doc.page_content for doc in documents])
            metadata = {
                "presentation": {
                    "slide_count": len(documents)
                }
            }
            return text_content, metadata
        except Exception as e:
            logger.error(f"Error extracting PowerPoint {file_path}: {e}", exc_info=True)
            return "", None
    
    async def _extract_html(self, file_path: str) -> str:
        """Extract text from HTML file."""
        try:
            loader = UnstructuredHTMLLoader(file_path)
            documents = loader.load()
            return "\n\n".join([doc.page_content for doc in documents])
        except Exception as e:
            logger.error(f"Error extracting HTML {file_path}: {e}")
            return ""
    
    async def _extract_markdown(self, file_path: str) -> str:
        """Extract text from Markdown file, preserving structure."""
        # Markdown extraction is required - no fallbacks
        logger.info(f"Extracting text from Markdown: {file_path}")
        
        loader = UnstructuredMarkdownLoader(file_path)
        documents = loader.load()
        
        if not documents:
            error_msg = f"No content extracted from Markdown: {file_path}"
            logger.error(error_msg)
            raise ValueError(error_msg)
        
        # Combine documents, preserving markdown structure
        # UnstructuredMarkdownLoader already handles markdown parsing
        text_content = "\n\n".join([doc.page_content for doc in documents])
        
        if not text_content or not text_content.strip():
            error_msg = f"Empty content extracted from Markdown: {file_path}"
            logger.error(error_msg)
            raise ValueError(error_msg)
        
        # Log extraction stats
        text_length = len(text_content)
        header_count = len(re.findall(r'^#{1,6}\s+', text_content, re.MULTILINE))
        code_block_count = len(re.findall(r'```', text_content))
        logger.info(f"Extracted {text_length} characters from Markdown: {file_path} "
                   f"(headers: {header_count}, code blocks: {code_block_count})")
        
        return text_content
    
    async def _extract_text_file(self, file_path: str) -> str:
        """Extract text from plain text file."""
        try:
            loader = TextLoader(file_path, encoding='utf-8')
            documents = loader.load()
            return "\n\n".join([doc.page_content for doc in documents])
        except Exception as e:
            logger.error(f"Error extracting text file {file_path}: {e}")
            # Fallback to direct file reading
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    return f.read()
            except Exception as e2:
                logger.error(f"Fallback text extraction failed for {file_path}: {e2}")
                return ""

    def _get_shape_text(self, shape) -> str:
        """Extract text from a pptx shape."""
        try:
            if not getattr(shape, "has_text_frame", False):
                return ""
            paragraphs = []
            for paragraph in shape.text_frame.paragraphs:
                runs = [run.text for run in paragraph.runs]
                text = "".join(runs).strip()
                if text:
                    paragraphs.append(text)
            return "\n".join(paragraphs).strip()
        except Exception:
            return ""
    
    async def split_text_with_metadata(
        self,
        text: str,
        chunk_size: int = 1000,
        chunk_overlap: int = 200,
        strategy: Optional[str] = None
    ) -> List[Dict[str, Any]]:
        """
        Split text into chunks with enhanced metadata.
        
        Args:
            text: Text content to split
            chunk_size: Maximum size of each chunk in characters
            chunk_overlap: Number of characters to overlap between chunks
            strategy: Chunking strategy ('semantic' or 'fixed')
            
        Returns:
            List of dictionaries with chunk content and metadata
        """
        chunks = await self.split_text(text, chunk_size, chunk_overlap, strategy)
        
        # Extract metadata for each chunk
        chunks_with_metadata = []
        current_pos = 0
        
        for i, chunk_content in enumerate(chunks):
            # Find position in original text
            start_pos = text.find(chunk_content, current_pos)
            if start_pos == -1:
                start_pos = current_pos
            end_pos = start_pos + len(chunk_content)
            current_pos = start_pos
            
            # Extract section title if available (look for headers before chunk)
            section_title = self._extract_section_title(text, start_pos)
            
            # Calculate semantic score (placeholder - could use embedding similarity)
            semantic_score = 1.0  # Default, could be calculated based on coherence
            
            chunk_metadata = {
                "content": chunk_content,
                "chunk_index": i,
                "start_pos": start_pos,
                "end_pos": end_pos,
                "section_title": section_title,
                "semantic_score": semantic_score,
                "paragraph_index": self._get_paragraph_index(text, start_pos)
            }
            
            chunks_with_metadata.append(chunk_metadata)
        
        return chunks_with_metadata
    
    def _extract_section_title(self, text: str, position: int) -> Optional[str]:
        """Extract section title (header) before the given position."""
        # Look for markdown headers or HTML headers before position
        before_text = text[:position]
        
        # Find last markdown header
        markdown_headers = re.findall(r'^#{1,6}\s+(.+)$', before_text, re.MULTILINE)
        if markdown_headers:
            return markdown_headers[-1].strip()
        
        # Find last HTML header
        html_headers = re.findall(r'<h[1-6][^>]*>(.+?)</h[1-6]>', before_text, re.IGNORECASE | re.DOTALL)
        if html_headers:
            # Clean HTML tags
            header_text = re.sub(r'<[^>]+>', '', html_headers[-1])
            return header_text.strip()
        
        return None
    
    def _get_paragraph_index(self, text: str, position: int) -> int:
        """Get paragraph index for the given position."""
        before_text = text[:position]
        paragraphs = re.split(r'\n\s*\n', before_text)
        return len([p for p in paragraphs if p.strip()])
    
    async def split_text(
        self,
        text: str,
        chunk_size: int = 1000,
        chunk_overlap: int = 200,
        strategy: Optional[str] = None
    ) -> List[str]:
        """
        Split text into chunks for processing.
        
        Args:
            text: Text content to split
            chunk_size: Maximum size of each chunk in characters (default: 1000)
            chunk_overlap: Number of characters to overlap between chunks (default: 200)
            strategy: Chunking strategy ('semantic' or 'fixed'), uses config if None
            
        Returns:
            List of text chunks, filtered to remove very short chunks (< 50 chars)
            
        Note:
            Chunks shorter than 50 characters are filtered out to avoid
            processing noise or incomplete sentences.
        """
        if not text or not text.strip():
            return []
        
        strategy = strategy or settings.RAG_CHUNKING_STRATEGY
        
        try:
            # Use semantic chunking if enabled and model available
            if strategy == "semantic" and self.semantic_model:
                return await self._semantic_chunk(text, chunk_size, chunk_overlap)
            else:
                # Use fixed-size chunking
                return await self._fixed_chunk(text, chunk_size, chunk_overlap)
            
        except Exception as e:
            logger.error(f"Error splitting text: {e}")
            return [text]  # Return original text as single chunk
    
    async def _fixed_chunk(
        self,
        text: str,
        chunk_size: int,
        chunk_overlap: int
    ) -> List[str]:
        """Split text using fixed-size chunking."""
        # Detect if text is markdown and use markdown-aware separators
        is_markdown = self._is_markdown(text)
        
        # Update splitter configuration if different from default
        if chunk_size != self.text_splitter._chunk_size or chunk_overlap != self.text_splitter._chunk_overlap:
            # Use markdown-aware separators if content is markdown
            # This helps preserve structure: headers, code blocks, lists
            separators = (
                ["\n\n## ", "\n\n# ", "\n\n```", "\n\n", "\n", " ", ""] 
                if is_markdown 
                else ["\n\n", "\n", " ", ""]
            )
            
            self.text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=chunk_size,
                chunk_overlap=chunk_overlap,
                length_function=len,
                separators=separators
            )
        
        # Clean the text
        cleaned_text = self._clean_text(text)
        
        # Split into chunks
        chunks = self.text_splitter.split_text(cleaned_text)
        
        # Filter out very short chunks
        min_chunk_length = 50
        filtered_chunks = [chunk for chunk in chunks if len(chunk.strip()) >= min_chunk_length]
        
        logger.debug(f"Split text into {len(filtered_chunks)} fixed-size chunks")
        return filtered_chunks
    
    async def _semantic_chunk(
        self,
        text: str,
        max_chunk_size: int,
        overlap: int
    ) -> List[str]:
        """
        Split text using semantic chunking based on topic boundaries.
        
        Args:
            text: Text to chunk
            max_chunk_size: Maximum chunk size in characters
            overlap: Overlap between chunks in characters
            
        Returns:
            List of semantically coherent chunks
        """
        if not self.semantic_model:
            logger.warning("Semantic model not available, falling back to fixed chunking")
            return await self._fixed_chunk(text, max_chunk_size, overlap)
        
        try:
            # Clean the text
            cleaned_text = self._clean_text(text)
            
            # For markdown, split by headers first to preserve structure
            is_markdown = self._is_markdown(cleaned_text)
            if is_markdown:
                # Split by markdown headers (## or #) to create sections
                # This preserves markdown structure better
                sections = re.split(r'\n(#{1,6}\s+)', cleaned_text)
                # Recombine headers with their content
                paragraphs = []
                for i in range(0, len(sections) - 1, 2):
                    if i + 1 < len(sections):
                        paragraphs.append(sections[i] + sections[i + 1])
                    else:
                        paragraphs.append(sections[i])
                # Also split by double newlines within sections
                all_paragraphs = []
                for section in paragraphs:
                    all_paragraphs.extend(re.split(r'\n\s*\n', section))
            else:
                # Split by paragraphs first (natural boundaries)
                paragraphs = re.split(r'\n\s*\n', cleaned_text)
                all_paragraphs = paragraphs
            
            # Group paragraphs into semantic chunks
            chunks = []
            current_chunk = []
            current_size = 0
            
            for para in all_paragraphs:
                para = para.strip()
                if not para:
                    continue
                
                para_size = len(para)
                
                # If adding this paragraph would exceed max size, finalize current chunk
                if current_size + para_size > max_chunk_size and current_chunk:
                    chunk_text = "\n\n".join(current_chunk)
                    if len(chunk_text.strip()) >= 50:  # Minimum chunk size
                        chunks.append(chunk_text)
                    
                    # Start new chunk with overlap
                    if overlap > 0 and current_chunk:
                        # Take last part of previous chunk for overlap
                        overlap_text = current_chunk[-1][-overlap:] if current_chunk else ""
                        current_chunk = [overlap_text, para] if overlap_text else [para]
                        current_size = len("\n\n".join(current_chunk))
                    else:
                        current_chunk = [para]
                        current_size = para_size
                else:
                    current_chunk.append(para)
                    current_size += para_size + 2  # +2 for "\n\n"
            
            # Add final chunk
            if current_chunk:
                chunk_text = "\n\n".join(current_chunk)
                if len(chunk_text.strip()) >= 50:
                    chunks.append(chunk_text)
            
            # If we have very few chunks, try sentence-level semantic splitting
            if len(chunks) < 2 and len(cleaned_text) > max_chunk_size:
                chunks = await self._semantic_chunk_by_sentences(cleaned_text, max_chunk_size, overlap)
            
            logger.debug(f"Split text into {len(chunks)} semantic chunks")
            return chunks
            
        except Exception as e:
            logger.error(f"Error in semantic chunking: {e}, falling back to fixed chunking")
            return await self._fixed_chunk(text, max_chunk_size, overlap)
    
    async def _semantic_chunk_by_sentences(
        self,
        text: str,
        max_chunk_size: int,
        overlap: int
    ) -> List[str]:
        """Split text by sentences using semantic similarity."""
        # Split into sentences
        sentences = re.split(r'[.!?]+\s+', text)
        sentences = [s.strip() for s in sentences if s.strip()]
        
        if not sentences:
            return [text]
        
        # Generate embeddings for sentences
        embeddings = self.semantic_model.encode(sentences, show_progress_bar=False)
        
        chunks = []
        current_chunk = []
        current_size = 0
        last_embedding = None
        
        for i, sentence in enumerate(sentences):
            sentence_size = len(sentence)
            embedding = embeddings[i]
            
            # Calculate similarity to previous sentence
            similarity = 1.0
            if last_embedding is not None:
                similarity = np.dot(embedding, last_embedding) / (
                    np.linalg.norm(embedding) * np.linalg.norm(last_embedding)
                )
            
            # If similarity is low (topic change) or size limit reached, start new chunk
            if (similarity < 0.7 and current_chunk) or (current_size + sentence_size > max_chunk_size and current_chunk):
                chunk_text = " ".join(current_chunk)
                if len(chunk_text.strip()) >= 50:
                    chunks.append(chunk_text)
                current_chunk = [sentence]
                current_size = sentence_size
            else:
                current_chunk.append(sentence)
                current_size += sentence_size + 1  # +1 for space
            
            last_embedding = embedding
        
        # Add final chunk
        if current_chunk:
            chunk_text = " ".join(current_chunk)
            if len(chunk_text.strip()) >= 50:
                chunks.append(chunk_text)
        
        return chunks if chunks else [text]
    
    def _is_markdown(self, text: str) -> bool:
        """
        Detect if text content is markdown format.
        
        Args:
            text: Text content to check
            
        Returns:
            True if text appears to be markdown, False otherwise
        """
        if not text:
            return False
        
        # Check for common markdown patterns
        markdown_patterns = [
            r'^#{1,6}\s+',  # Headers (# Header)
            r'\*\*.*?\*\*',  # Bold (**text**)
            r'\*.*?\*',  # Italic (*text*)
            r'\[.*?\]\(.*?\)',  # Links [text](url)
            r'```',  # Code blocks
            r'^\s*[-*+]\s+',  # Unordered lists
            r'^\s*\d+\.\s+',  # Ordered lists
            r'^\s*>\s+',  # Blockquotes
        ]
        
        # Count markdown patterns found
        pattern_count = sum(1 for pattern in markdown_patterns if re.search(pattern, text, re.MULTILINE))
        
        # Consider it markdown if at least 2 patterns are found
        # This helps avoid false positives on plain text with occasional markdown-like characters
        return pattern_count >= 2
    
    def _clean_text(self, text: str) -> str:
        """Clean and normalize text content."""
        if not text:
            return ""
        
        # Remove excessive whitespace
        text = re.sub(r'\s+', ' ', text)
        
        # Remove special characters that might cause issues
        text = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]', '', text)
        
        # Normalize line breaks
        text = re.sub(r'\r\n|\r', '\n', text)
        
        # Remove excessive blank lines
        text = re.sub(r'\n\s*\n\s*\n', '\n\n', text)
        
        return text.strip()
    
    async def extract_metadata(self, file_path: str) -> Dict[str, Any]:
        """
        Extract metadata from a file.
        
        Args:
            file_path: Path to the file
            
        Returns:
            Dictionary containing file metadata:
            - file_path: Full path to the file
            - file_name: Base name of the file
            - file_extension: File extension (lowercase)
            - file_size: File size in bytes
            - created_time: File creation timestamp (if available)
            - modified_time: File modification timestamp (if available)
        """
        metadata = {
            "file_path": file_path,
            "file_name": os.path.basename(file_path),
            "file_extension": os.path.splitext(file_path)[1].lower(),
            "file_size": 0
        }
        
        try:
            if os.path.exists(file_path):
                stat = os.stat(file_path)
                metadata.update({
                    "file_size": stat.st_size,
                    "created_time": stat.st_ctime,
                    "modified_time": stat.st_mtime
                })
        except Exception as e:
            logger.error(f"Error extracting metadata from {file_path}: {e}")
        
        return metadata
    
    def get_supported_formats(self) -> List[str]:
        """Get list of supported file formats."""
        return [
            'pdf', 'txt', 'docx', 'doc', 'pptx', 'ppt', 'html', 'htm', 'md', 'markdown'
        ]
    
    def is_supported_format(self, file_path: str) -> bool:
        """Check if file format is supported."""
        file_extension = os.path.splitext(file_path)[1].lower().lstrip('.')
        return file_extension in self.get_supported_formats()
