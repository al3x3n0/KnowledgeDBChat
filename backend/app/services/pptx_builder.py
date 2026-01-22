"""
PowerPoint presentation builder service.

Creates PPTX files from presentation outlines using python-pptx.
"""

from io import BytesIO
from typing import Dict, Optional, Any
import os
from loguru import logger

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from app.schemas.presentation import PresentationOutline, SlideContent


# Mapping from our slide types to standard PowerPoint layout indices
# These are the default indices for layouts in a standard PPTX template
LAYOUT_MAPPING = {
    "title": 0,       # Title Slide
    "content": 1,     # Title and Content
    "two_column": 3,  # Two Content
    "diagram": 6,     # Blank (for custom placement)
    "summary": 1,     # Title and Content
}


def hex_to_rgb(hex_color: str) -> RgbColor:
    """Convert hex color string to RgbColor."""
    hex_color = hex_color.lstrip('#')
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return RgbColor(r, g, b)


class PPTXBuilder:
    """
    Builds PowerPoint presentations from outline data.

    Supports multiple built-in styles, custom themes, and various slide types.
    """

    # Slide dimensions (16:9 aspect ratio)
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)

    # Built-in style definitions
    STYLES = {
        "professional": {
            "title_color": RgbColor(0x1a, 0x36, 0x5d),      # Dark blue
            "accent_color": RgbColor(0x2e, 0x86, 0xab),      # Light blue
            "text_color": RgbColor(0x33, 0x33, 0x33),        # Dark gray
            "bg_color": RgbColor(0xff, 0xff, 0xff),          # White
            "title_font": "Calibri",
            "body_font": "Calibri",
            "title_size": Pt(44),
            "subtitle_size": Pt(24),
            "heading_size": Pt(36),
            "body_size": Pt(20),
            "bullet_size": Pt(18),
        },
        "casual": {
            "title_color": RgbColor(0x4a, 0x90, 0xd9),      # Sky blue
            "accent_color": RgbColor(0xff, 0x6b, 0x6b),      # Coral
            "text_color": RgbColor(0x2d, 0x3a, 0x4a),        # Dark slate
            "bg_color": RgbColor(0xf8, 0xf9, 0xfa),          # Light gray
            "title_font": "Arial Rounded MT Bold",
            "body_font": "Arial",
            "title_size": Pt(48),
            "subtitle_size": Pt(26),
            "heading_size": Pt(38),
            "body_size": Pt(22),
            "bullet_size": Pt(20),
        },
        "technical": {
            "title_color": RgbColor(0x00, 0x7a, 0xcc),      # Tech blue
            "accent_color": RgbColor(0x28, 0xa7, 0x45),      # Green
            "text_color": RgbColor(0x24, 0x29, 0x2e),        # GitHub dark
            "bg_color": RgbColor(0xff, 0xff, 0xff),          # White
            "title_font": "Consolas",
            "body_font": "Segoe UI",
            "title_size": Pt(40),
            "subtitle_size": Pt(22),
            "heading_size": Pt(32),
            "body_size": Pt(18),
            "bullet_size": Pt(16),
        },
        "modern": {
            "title_color": RgbColor(0x2c, 0x3e, 0x50),      # Midnight blue
            "accent_color": RgbColor(0xe7, 0x4c, 0x3c),      # Red
            "text_color": RgbColor(0x34, 0x49, 0x5e),        # Wet asphalt
            "bg_color": RgbColor(0xec, 0xf0, 0xf1),          # Clouds
            "title_font": "Segoe UI",
            "body_font": "Segoe UI",
            "title_size": Pt(46),
            "subtitle_size": Pt(24),
            "heading_size": Pt(36),
            "body_size": Pt(20),
            "bullet_size": Pt(18),
        },
        "minimal": {
            "title_color": RgbColor(0x00, 0x00, 0x00),      # Black
            "accent_color": RgbColor(0x95, 0xa5, 0xa6),      # Gray
            "text_color": RgbColor(0x2c, 0x3e, 0x50),        # Dark blue-gray
            "bg_color": RgbColor(0xff, 0xff, 0xff),          # White
            "title_font": "Helvetica",
            "body_font": "Helvetica",
            "title_size": Pt(42),
            "subtitle_size": Pt(22),
            "heading_size": Pt(34),
            "body_size": Pt(18),
            "bullet_size": Pt(16),
        },
        "corporate": {
            "title_color": RgbColor(0x00, 0x3d, 0x7a),      # Corporate blue
            "accent_color": RgbColor(0xf5, 0xa6, 0x23),      # Orange
            "text_color": RgbColor(0x33, 0x33, 0x33),        # Dark gray
            "bg_color": RgbColor(0xff, 0xff, 0xff),          # White
            "title_font": "Arial",
            "body_font": "Arial",
            "title_size": Pt(44),
            "subtitle_size": Pt(24),
            "heading_size": Pt(36),
            "body_size": Pt(20),
            "bullet_size": Pt(18),
        },
        "creative": {
            "title_color": RgbColor(0x9b, 0x59, 0xb6),      # Purple
            "accent_color": RgbColor(0x1a, 0xbc, 0x9c),      # Turquoise
            "text_color": RgbColor(0x2c, 0x3e, 0x50),        # Dark
            "bg_color": RgbColor(0xfd, 0xfb, 0xf7),          # Off-white
            "title_font": "Georgia",
            "body_font": "Verdana",
            "title_size": Pt(48),
            "subtitle_size": Pt(26),
            "heading_size": Pt(38),
            "body_size": Pt(20),
            "bullet_size": Pt(18),
        },
        "dark": {
            "title_color": RgbColor(0xff, 0xff, 0xff),      # White
            "accent_color": RgbColor(0x3d, 0xb9, 0xd3),      # Cyan
            "text_color": RgbColor(0xe0, 0xe0, 0xe0),        # Light gray
            "bg_color": RgbColor(0x1e, 0x1e, 0x2e),          # Dark background
            "title_font": "Segoe UI",
            "body_font": "Segoe UI",
            "title_size": Pt(44),
            "subtitle_size": Pt(24),
            "heading_size": Pt(36),
            "body_size": Pt(20),
            "bullet_size": Pt(18),
        },
    }

    def __init__(
        self,
        style: str = "professional",
        custom_theme: Optional[Dict[str, Any]] = None,
        template_path: Optional[str] = None
    ):
        """
        Initialize the builder with a style, custom theme, or template file.

        Args:
            style: Built-in style name (used if no custom_theme or template_path)
            custom_theme: Custom theme configuration dict
            template_path: Path to a PPTX template file to use as base
        """
        self.style = style
        self.template_path = template_path
        self._use_template = False

        # Check if template file exists and is valid
        if template_path and os.path.exists(template_path):
            try:
                # Verify we can open the template
                test_prs = Presentation(template_path)
                self._use_template = True
                self._template_layout_count = len(test_prs.slide_layouts)
                logger.info(f"Using PPTX template with {self._template_layout_count} layouts")
            except Exception as e:
                logger.warning(f"Failed to load template, falling back to default: {e}")
                self._use_template = False

        # Parse theme configuration
        if custom_theme:
            self.style_config = self._parse_custom_theme(custom_theme)
        else:
            self.style_config = self.STYLES.get(style, self.STYLES["professional"])

    def _parse_custom_theme(self, theme: Dict[str, Any]) -> Dict[str, Any]:
        """
        Parse custom theme configuration into style config.

        Args:
            theme: Custom theme dict with colors, fonts, sizes sections

        Returns:
            Style config dict compatible with built-in styles
        """
        # Start with professional defaults
        config = dict(self.STYLES["professional"])

        # Parse colors
        colors = theme.get("colors", {})
        if "title_color" in colors:
            config["title_color"] = hex_to_rgb(colors["title_color"])
        if "accent_color" in colors:
            config["accent_color"] = hex_to_rgb(colors["accent_color"])
        if "text_color" in colors:
            config["text_color"] = hex_to_rgb(colors["text_color"])
        if "bg_color" in colors:
            config["bg_color"] = hex_to_rgb(colors["bg_color"])

        # Parse fonts
        fonts = theme.get("fonts", {})
        if "title_font" in fonts:
            config["title_font"] = fonts["title_font"]
        if "body_font" in fonts:
            config["body_font"] = fonts["body_font"]

        # Parse sizes
        sizes = theme.get("sizes", {})
        if "title_size" in sizes:
            config["title_size"] = Pt(sizes["title_size"])
        if "subtitle_size" in sizes:
            config["subtitle_size"] = Pt(sizes["subtitle_size"])
        if "heading_size" in sizes:
            config["heading_size"] = Pt(sizes["heading_size"])
        if "body_size" in sizes:
            config["body_size"] = Pt(sizes["body_size"])
        if "bullet_size" in sizes:
            config["bullet_size"] = Pt(sizes["bullet_size"])

        return config

    @classmethod
    def get_available_styles(cls) -> Dict[str, Dict[str, str]]:
        """
        Get list of available built-in styles with descriptions.

        Returns:
            Dict mapping style name to description and color preview
        """
        return {
            "professional": {
                "name": "Professional",
                "description": "Clean, corporate look with dark blue accents",
                "primary_color": "#1a365d",
                "accent_color": "#2e86ab",
            },
            "casual": {
                "name": "Casual",
                "description": "Friendly and approachable with warm colors",
                "primary_color": "#4a90d9",
                "accent_color": "#ff6b6b",
            },
            "technical": {
                "name": "Technical",
                "description": "Developer-focused with monospace fonts",
                "primary_color": "#007acc",
                "accent_color": "#28a745",
            },
            "modern": {
                "name": "Modern",
                "description": "Contemporary design with bold contrasts",
                "primary_color": "#2c3e50",
                "accent_color": "#e74c3c",
            },
            "minimal": {
                "name": "Minimal",
                "description": "Simple and clean with subtle grays",
                "primary_color": "#000000",
                "accent_color": "#95a5a6",
            },
            "corporate": {
                "name": "Corporate",
                "description": "Traditional business style with orange accents",
                "primary_color": "#003d7a",
                "accent_color": "#f5a623",
            },
            "creative": {
                "name": "Creative",
                "description": "Artistic with purple and turquoise",
                "primary_color": "#9b59b6",
                "accent_color": "#1abc9c",
            },
            "dark": {
                "name": "Dark",
                "description": "Dark theme for low-light presentations",
                "primary_color": "#ffffff",
                "accent_color": "#3db9d3",
            },
        }

    def build(
        self,
        outline: PresentationOutline,
        diagrams: Optional[Dict[int, bytes]] = None
    ) -> bytes:
        """
        Build a PowerPoint presentation from an outline.

        Args:
            outline: The presentation outline with slide content
            diagrams: Optional dict mapping slide_number to PNG bytes

        Returns:
            PPTX file as bytes
        """
        diagrams = diagrams or {}

        # Create or load presentation
        if self._use_template and self.template_path:
            logger.info(f"Loading PPTX template from {self.template_path}")
            prs = Presentation(self.template_path)
            # Remove any existing slides from template (we'll add our own content)
            # Keep slide masters and layouts intact
            while len(prs.slides) > 0:
                rId = prs.slides._sldIdLst[0].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[0]
        else:
            prs = Presentation()
            prs.slide_width = self.SLIDE_WIDTH
            prs.slide_height = self.SLIDE_HEIGHT

        # Build each slide
        for slide_content in outline.slides:
            slide_type = slide_content.slide_type

            if self._use_template:
                # Use template layouts
                self._add_slide_from_template(prs, slide_content, diagrams)
            else:
                # Use custom slide builders
                if slide_type == "title":
                    self._add_title_slide(prs, slide_content)
                elif slide_type == "content":
                    self._add_content_slide(prs, slide_content)
                elif slide_type == "diagram":
                    diagram_png = diagrams.get(slide_content.slide_number)
                    self._add_diagram_slide(prs, slide_content, diagram_png)
                elif slide_type == "summary":
                    self._add_summary_slide(prs, slide_content)
                elif slide_type == "two_column":
                    self._add_two_column_slide(prs, slide_content)
                else:
                    # Default to content slide
                    self._add_content_slide(prs, slide_content)

        # Save to bytes
        buffer = BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()

    def _get_template_layout(self, prs: Presentation, slide_type: str):
        """
        Get the appropriate layout from the template for a given slide type.

        Args:
            prs: The presentation object
            slide_type: Type of slide (title, content, diagram, etc.)

        Returns:
            Slide layout from the template
        """
        layout_idx = LAYOUT_MAPPING.get(slide_type, 1)  # Default to content

        try:
            return prs.slide_layouts[layout_idx]
        except IndexError:
            # Fallback to first available layout
            logger.warning(f"Layout index {layout_idx} not found, using fallback")
            return prs.slide_layouts[0] if len(prs.slide_layouts) > 0 else None

    def _add_slide_from_template(
        self,
        prs: Presentation,
        content: SlideContent,
        diagrams: Dict[int, bytes]
    ):
        """
        Add a slide using the template's layouts.

        This method uses the template's built-in layouts and placeholders
        to maintain the template's design while adding our content.
        """
        slide_type = content.slide_type
        layout = self._get_template_layout(prs, slide_type)

        if layout is None:
            logger.error("No layout available in template")
            return

        slide = prs.slides.add_slide(layout)

        # Try to populate placeholders
        for shape in slide.placeholders:
            ph_type = shape.placeholder_format.type
            # Placeholder types: TITLE=1, BODY=2, CENTER_TITLE=3, SUBTITLE=4, etc.

            if ph_type in (1, 3):  # TITLE or CENTER_TITLE
                shape.text = content.title
            elif ph_type == 4:  # SUBTITLE
                if content.subtitle:
                    shape.text = content.subtitle
                elif content.content and slide_type == "title":
                    shape.text = content.content[0] if content.content else ""
            elif ph_type == 2:  # BODY
                if content.content:
                    tf = shape.text_frame
                    for i, point in enumerate(content.content):
                        if i == 0:
                            tf.paragraphs[0].text = point
                        else:
                            p = tf.add_paragraph()
                            p.text = point
                            p.level = 0

        # Handle diagram slides specially
        if slide_type == "diagram":
            diagram_png = diagrams.get(content.slide_number)
            if diagram_png:
                # Add diagram image to the slide
                image_stream = BytesIO(diagram_png)
                # Position in center of slide
                img_left = Inches(1.5)
                img_top = Inches(1.8)
                img_width = Inches(10)
                slide.shapes.add_picture(image_stream, img_left, img_top, width=img_width)

        # Add speaker notes
        if content.notes:
            slide.notes_slide.notes_text_frame.text = content.notes

    def _add_title_slide(self, prs: Presentation, content: SlideContent):
        """Add a title slide."""
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)

        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.paragraphs[0].text = content.title
        title_frame.paragraphs[0].font.size = self.style_config["title_size"]
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.style_config["title_color"]
        title_frame.paragraphs[0].font.name = self.style_config["title_font"]
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Add subtitle if present
        subtitle = content.subtitle or (content.content[0] if content.content else None)
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8)
            )
            sub_frame = subtitle_box.text_frame
            sub_frame.paragraphs[0].text = subtitle
            sub_frame.paragraphs[0].font.size = self.style_config["subtitle_size"]
            sub_frame.paragraphs[0].font.color.rgb = self.style_config["text_color"]
            sub_frame.paragraphs[0].font.name = self.style_config["body_font"]
            sub_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Add decorative line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(4), Inches(4.0), Inches(5.333), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.style_config["accent_color"]
        line.line.fill.background()

        # Add speaker notes
        if content.notes:
            slide.notes_slide.notes_text_frame.text = content.notes

    def _add_content_slide(self, prs: Presentation, content: SlideContent):
        """Add a content slide with bullet points."""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Add heading
        heading_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(1)
        )
        heading_frame = heading_box.text_frame
        heading_frame.paragraphs[0].text = content.title
        heading_frame.paragraphs[0].font.size = self.style_config["heading_size"]
        heading_frame.paragraphs[0].font.bold = True
        heading_frame.paragraphs[0].font.color.rgb = self.style_config["title_color"]
        heading_frame.paragraphs[0].font.name = self.style_config["title_font"]

        # Add underline accent
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1.2), Inches(2), Inches(0.04)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.style_config["accent_color"]
        line.line.fill.background()

        # Add bullet points
        if content.content:
            content_box = slide.shapes.add_textbox(
                Inches(0.7), Inches(1.5), Inches(11.5), Inches(5.5)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True

            for i, point in enumerate(content.content):
                if i == 0:
                    p = content_frame.paragraphs[0]
                else:
                    p = content_frame.add_paragraph()

                p.text = f"• {point}"
                p.font.size = self.style_config["bullet_size"]
                p.font.color.rgb = self.style_config["text_color"]
                p.font.name = self.style_config["body_font"]
                p.space_after = Pt(12)
                p.level = 0

        # Add speaker notes
        if content.notes:
            slide.notes_slide.notes_text_frame.text = content.notes

    def _add_diagram_slide(
        self,
        prs: Presentation,
        content: SlideContent,
        diagram_png: Optional[bytes] = None
    ):
        """Add a slide with a diagram image."""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Add heading
        heading_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8)
        )
        heading_frame = heading_box.text_frame
        heading_frame.paragraphs[0].text = content.title
        heading_frame.paragraphs[0].font.size = self.style_config["heading_size"]
        heading_frame.paragraphs[0].font.bold = True
        heading_frame.paragraphs[0].font.color.rgb = self.style_config["title_color"]
        heading_frame.paragraphs[0].font.name = self.style_config["title_font"]

        if diagram_png:
            # Add the diagram image
            image_stream = BytesIO(diagram_png)
            # Center the image with reasonable size
            img_left = Inches(1.5)
            img_top = Inches(1.3)
            img_width = Inches(10)

            slide.shapes.add_picture(
                image_stream,
                img_left,
                img_top,
                width=img_width
            )
        else:
            # Placeholder if no diagram available
            placeholder_box = slide.shapes.add_textbox(
                Inches(2), Inches(3), Inches(9), Inches(1)
            )
            placeholder_frame = placeholder_box.text_frame
            placeholder_frame.paragraphs[0].text = "[Diagram could not be rendered]"
            placeholder_frame.paragraphs[0].font.size = Pt(18)
            placeholder_frame.paragraphs[0].font.italic = True
            placeholder_frame.paragraphs[0].font.color.rgb = RgbColor(0x99, 0x99, 0x99)
            placeholder_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Show the diagram description if available
            if content.diagram_description:
                desc_box = slide.shapes.add_textbox(
                    Inches(2), Inches(4), Inches(9), Inches(2)
                )
                desc_frame = desc_box.text_frame
                desc_frame.paragraphs[0].text = f"Description: {content.diagram_description}"
                desc_frame.paragraphs[0].font.size = Pt(14)
                desc_frame.paragraphs[0].font.color.rgb = RgbColor(0x66, 0x66, 0x66)
                desc_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Add speaker notes
        if content.notes:
            slide.notes_slide.notes_text_frame.text = content.notes

    def _add_summary_slide(self, prs: Presentation, content: SlideContent):
        """Add a summary/conclusion slide."""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Add heading with different styling for summary
        heading_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(1)
        )
        heading_frame = heading_box.text_frame
        heading_frame.paragraphs[0].text = content.title
        heading_frame.paragraphs[0].font.size = self.style_config["heading_size"]
        heading_frame.paragraphs[0].font.bold = True
        heading_frame.paragraphs[0].font.color.rgb = self.style_config["title_color"]
        heading_frame.paragraphs[0].font.name = self.style_config["title_font"]

        # Add accent bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(1.1), Inches(13.333), Inches(0.08)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.style_config["accent_color"]
        bar.line.fill.background()

        # Add key points
        if content.content:
            content_box = slide.shapes.add_textbox(
                Inches(0.7), Inches(1.5), Inches(11.5), Inches(5.5)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True

            for i, point in enumerate(content.content):
                if i == 0:
                    p = content_frame.paragraphs[0]
                else:
                    p = content_frame.add_paragraph()

                # Use checkmark for summary points
                p.text = f"✓ {point}"
                p.font.size = self.style_config["body_size"]
                p.font.color.rgb = self.style_config["text_color"]
                p.font.name = self.style_config["body_font"]
                p.space_after = Pt(16)

        # Add speaker notes
        if content.notes:
            slide.notes_slide.notes_text_frame.text = content.notes

    def _add_two_column_slide(self, prs: Presentation, content: SlideContent):
        """Add a two-column layout slide."""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Add heading
        heading_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8)
        )
        heading_frame = heading_box.text_frame
        heading_frame.paragraphs[0].text = content.title
        heading_frame.paragraphs[0].font.size = self.style_config["heading_size"]
        heading_frame.paragraphs[0].font.bold = True
        heading_frame.paragraphs[0].font.color.rgb = self.style_config["title_color"]

        # Split content into two columns
        if content.content:
            mid = len(content.content) // 2

            # Left column
            left_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.3), Inches(5.8), Inches(5.5)
            )
            left_frame = left_box.text_frame
            left_frame.word_wrap = True

            for i, point in enumerate(content.content[:mid]):
                if i == 0:
                    p = left_frame.paragraphs[0]
                else:
                    p = left_frame.add_paragraph()
                p.text = f"• {point}"
                p.font.size = self.style_config["bullet_size"]
                p.font.color.rgb = self.style_config["text_color"]
                p.space_after = Pt(10)

            # Right column
            right_box = slide.shapes.add_textbox(
                Inches(6.8), Inches(1.3), Inches(5.8), Inches(5.5)
            )
            right_frame = right_box.text_frame
            right_frame.word_wrap = True

            for i, point in enumerate(content.content[mid:]):
                if i == 0:
                    p = right_frame.paragraphs[0]
                else:
                    p = right_frame.add_paragraph()
                p.text = f"• {point}"
                p.font.size = self.style_config["bullet_size"]
                p.font.color.rgb = self.style_config["text_color"]
                p.space_after = Pt(10)

        # Add speaker notes
        if content.notes:
            slide.notes_slide.notes_text_frame.text = content.notes
